from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建新的演示文稿
prs = Presentation()

# 设置幻灯片尺寸为16:9
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# 定义颜色
COLOR_PRIMARY = RGBColor(44, 62, 80)      # 深蓝
COLOR_SECONDARY = RGBColor(52, 152, 219)  # 浅蓝
COLOR_ACCENT = RGBColor(231, 76, 60)      # 红色
COLOR_WARNING = RGBColor(243, 156, 18)    # 橙色
COLOR_SUCCESS = RGBColor(46, 204, 113)    # 绿色
COLOR_TEXT = RGBColor(52, 73, 94)         # 深灰
COLOR_BG = RGBColor(236, 240, 241)        # 浅灰

# ========== 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
title = slide.shapes.add_textbox(Inches(2.5), Inches(1.5), Inches(5), Inches(1))
title_frame = title.text_frame
title_frame.text = "求职常见陷阱类型"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(44)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_PRIMARY
title_p.alignment = PP_ALIGN.CENTER

subtitle = slide.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(5), Inches(0.5))
subtitle_frame = subtitle.text_frame
subtitle_frame.text = "案例·特征·规避方法"
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.font.size = Pt(20)
subtitle_p.font.color.rgb = COLOR_SECONDARY
subtitle_p.alignment = PP_ALIGN.CENTER

footer = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(5), Inches(0.3))
footer_frame = footer.text_frame
footer_frame.text = "求职安全指南"
footer_p = footer_frame.paragraphs[0]
footer_p.font.size = Pt(14)
footer_p.font.color.rgb = RGBColor(149, 165, 166)
footer_p.alignment = PP_ALIGN.CENTER

# ========== 目录页 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "目录"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(32)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_PRIMARY

traps = [
    "虚假招聘陷阱",
    "培训贷陷阱",
    "收费面试陷阱",
    "高薪诱惑陷阱",
    "远程工作骗局",
    "实习转正骗局",
    "信息泄露陷阱",
    "加班文化隐性陷阱"
]

left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(4), Inches(3.5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4), Inches(3.5))
right_frame = right_box.text_frame
right_frame.word_wrap = True

for i, trap in enumerate(traps):
    if i < 4:
        p = left_frame.add_paragraph()
        p.text = f"{i+1}. {trap}"
    else:
        p = right_frame.add_paragraph()
        p.text = f"{i+1}. {trap}"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(12)

# ========== 简洁介绍页 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "求职陷阱的普遍性"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(32)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_PRIMARY

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "据调查，超过60%的求职者曾遭遇过不同类型的求职陷阱",
    "",
    "常见损失：",
    "• 经济损失：培训费、保证金、体检费等",
    "• 时间损失：无效面试、虚假培训",
    "• 信息损失：个人信息泄露、隐私被贩卖",
    "• 心理损失：信心受挫、职业规划被打乱",
    "",
    "本PPT将帮助你识别并规避这些求职陷阱"
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT
    if "•" in line:
        p.space_before = Pt(6)
        p.space_after = Pt(3)
    elif line == "":
        p.space_after = Pt(12)
    else:
        p.space_after = Pt(6)

# ========== 陷阱1: 虚假招聘 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "01 虚假招聘陷阱"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(28)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_ACCENT

# 标题背景
bg_box = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.15), Inches(0.5))
bg_box.fill.solid()
bg_box.fill.fore_color.rgb = COLOR_ACCENT
bg_box.line.fill.background()

# 案例区域
case_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(1.5))
case_frame = case_box.text_frame
case_frame.word_wrap = True

case_header = case_frame.paragraphs[0]
case_header.text = "典型案例"
case_header.font.size = Pt(18)
case_header.font.bold = True
case_header.font.color.rgb = COLOR_ACCENT
case_header.space_after = Pt(8)

case_content = case_frame.add_paragraph()
case_content.text = "某公司发布大量高薪职位，面试通过后要求先缴纳500元'入职保证金'，承诺3个月后退还，但公司随后消失或以各种理由拒绝退款。"
case_content.font.size = Pt(13)
case_content.font.color.rgb = COLOR_TEXT

# 特征区域
feature_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(1.5))
feature_frame = feature_box.text_frame
feature_frame.word_wrap = True

feature_header = feature_frame.paragraphs[0]
feature_header.text = "识别特征"
feature_header.font.size = Pt(18)
feature_header.font.bold = True
feature_header.font.color.rgb = COLOR_WARNING
feature_header.space_after = Pt(8)

features = ["职位描述模糊", "要求缴纳各类费用", "公司信息不透明", "薪资远高于市场水平"]
for f in features:
    p = feature_frame.add_paragraph()
    p.text = f"• {f}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

# 规避方法区域
avoid_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.8))
avoid_frame = avoid_box.text_frame
avoid_frame.word_wrap = True

avoid_header = avoid_frame.paragraphs[0]
avoid_header.text = "规避方法"
avoid_header.font.size = Pt(18)
avoid_header.font.bold = True
avoid_header.font.color.rgb = COLOR_SUCCESS
avoid_header.space_after = Pt(8)

avoids = [
    "查验企业资质：通过天眼查、企查查核实公司信息",
    "拒绝任何形式的付费要求：正规招聘不收取费用",
    "实地考察：条件允许时，去公司实地看看",
    "多方验证：在多个平台搜索公司评价"
]

for a in avoids:
    p = avoid_frame.add_paragraph()
    p.text = f"✓ {a}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

# ========== 陷阱2: 培训贷 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "02 培训贷陷阱"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(28)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_ACCENT

bg_box = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.15), Inches(0.5))
bg_box.fill.solid()
bg_box.fill.fore_color.rgb = COLOR_ACCENT
bg_box.line.fill.background()

case_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(1.5))
case_frame = case_box.text_frame
case_frame.word_wrap = True

case_header = case_frame.paragraphs[0]
case_header.text = "典型案例"
case_header.font.size = Pt(18)
case_header.font.bold = True
case_header.font.color.rgb = COLOR_ACCENT
case_header.space_after = Pt(8)

case_content = case_frame.add_paragraph()
case_content.text = "招聘声称提供'零基础入行'机会，要求先参加1-3个月培训，培训费可通过贷款支付，承诺培训后包就业。培训质量差，承诺的高薪职位不存在，学员背负高额贷款。"
case_content.font.size = Pt(13)
case_content.font.color.rgb = COLOR_TEXT

feature_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(1.5))
feature_frame = feature_box.text_frame
feature_frame.word_wrap = True

feature_header = feature_frame.paragraphs[0]
feature_header.text = "识别特征"
feature_header.font.size = Pt(18)
feature_header.font.bold = True
feature_header.font.color.rgb = COLOR_WARNING
feature_header.space_after = Pt(8)

features = ["承诺'零基础快速入行'", "培训费可通过贷款支付", "承诺培训后包就业", "培训内容与职位不匹配"]
for f in features:
    p = feature_frame.add_paragraph()
    p.text = f"• {f}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

avoid_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.8))
avoid_frame = avoid_box.text_frame
avoid_frame.word_wrap = True

avoid_header = avoid_frame.paragraphs[0]
avoid_header.text = "规避方法"
avoid_header.font.size = Pt(18)
avoid_header.font.bold = True
avoid_header.font.color.rgb = COLOR_SUCCESS
avoid_header.space_after = Pt(8)

avoids = [
    "警惕快速入行承诺：专业能力需要时间积累",
    "拒绝贷款培训：正规公司不会让你贷款工作",
    "考察培训机构资质：查询教育部门备案信息",
    "先签合同再交费：仔细阅读合同条款"
]

for a in avoids:
    p = avoid_frame.add_paragraph()
    p.text = f"✓ {a}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

# ========== 陷阱3: 收费面试 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "03 收费面试陷阱"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(28)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_ACCENT

bg_box = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.15), Inches(0.5))
bg_box.fill.solid()
bg_box.fill.fore_color.rgb = COLOR_ACCENT
bg_box.line.fill.background()

case_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(1.5))
case_frame = case_box.text_frame
case_frame.word_wrap = True

case_header = case_frame.paragraphs[0]
case_header.text = "典型案例"
case_header.font.size = Pt(18)
case_header.font.bold = True
case_header.font.color.rgb = COLOR_ACCENT
case_header.space_after = Pt(8)

case_content = case_frame.add_paragraph()
case_content.text = "面试通过后，HR告知入职前需要缴纳'保证金'、'服装费'、'资料费'等，合计2000元。缴费后，公司以各种理由推迟入职，最终彻底失联。"
case_content.font.size = Pt(13)
case_content.font.color.rgb = COLOR_TEXT

feature_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(1.5))
feature_frame = feature_box.text_frame
feature_frame.word_wrap = True

feature_header = feature_frame.paragraphs[0]
feature_header.text = "识别特征"
feature_header.font.size = Pt(18)
feature_header.font.bold = True
feature_header.font.color.rgb = COLOR_WARNING
feature_header.space_after = Pt(8)

features = ["入职前要求缴费", "费用种类繁多", "缴费后无正规票据", "公司拖延入职时间"]
for f in features:
    p = feature_frame.add_paragraph()
    p.text = f"• {f}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

avoid_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.8))
avoid_frame = avoid_box.text_frame
avoid_frame.word_wrap = True

avoid_header = avoid_frame.paragraphs[0]
avoid_header.text = "规避方法"
avoid_header.font.size = Pt(18)
avoid_header.font.bold = True
avoid_header.font.color.rgb = COLOR_SUCCESS
avoid_header.space_after = Pt(8)

avoids = [
    "坚守底线：正规公司不收取任何入职费用",
    "索要正规票据：如必须付费，要求开具发票",
    "保留付款凭证：截图、录音、合同等",
    "及时举报：向劳动监察部门投诉"
]

for a in avoids:
    p = avoid_frame.add_paragraph()
    p.text = f"✓ {a}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

# ========== 陷阱4: 高薪诱惑 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "04 高薪诱惑陷阱"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(28)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_ACCENT

bg_box = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.15), Inches(0.5))
bg_box.fill.solid()
bg_box.fill.fore_color.rgb = COLOR_ACCENT
bg_box.line.fill.background()

case_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(1.5))
case_frame = case_box.text_frame
case_frame.word_wrap = True

case_header = case_frame.paragraphs[0]
case_header.text = "典型案例"
case_header.font.size = Pt(18)
case_header.font.bold = True
case_header.font.color.rgb = COLOR_ACCENT
case_header.space_after = Pt(8)

case_content = case_frame.add_paragraph()
case_content.text = "某职位标注'月薪2万+，包吃住，周末双休'，实际工作内容为网络营销，底薪只有3000元，其余全靠业绩提成，工作强度极大，离职率极高。"
case_content.font.size = Pt(13)
case_content.font.color.rgb = COLOR_TEXT

feature_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(1.5))
feature_frame = feature_box.text_frame
feature_frame.word_wrap = True

feature_header = feature_frame.paragraphs[0]
feature_header.text = "识别特征"
feature_header.font.size = Pt(18)
feature_header.font.bold = True
feature_header.font.color.rgb = COLOR_WARNING
feature_header.space_after = Pt(8)

features = ["薪资远高于行业水平", "福利承诺过于完美", "职位描述模糊", "面试过程草率"]
for f in features:
    p = feature_frame.add_paragraph()
    p.text = f"• {f}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

avoid_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.8))
avoid_frame = avoid_box.text_frame
avoid_frame.word_wrap = True

avoid_header = avoid_frame.paragraphs[0]
avoid_header.text = "规避方法"
avoid_header.font.size = Pt(18)
avoid_header.font.bold = True
avoid_header.font.color.rgb = COLOR_SUCCESS
avoid_header.space_after = Pt(8)

avoids = [
    "调研薪资水平：了解行业正常薪资范围",
    "确认薪资结构：区分底薪、提成、绩效",
    "要求劳动合同明确：薪资数额写入合同",
    "警惕销售类岗位：高薪多依赖提成"
]

for a in avoids:
    p = avoid_frame.add_paragraph()
    p.text = f"✓ {a}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

# ========== 陷阱5: 远程工作骗局 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "05 远程工作骗局"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(28)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_ACCENT

bg_box = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.15), Inches(0.5))
bg_box.fill.solid()
bg_box.fill.fore_color.rgb = COLOR_ACCENT
bg_box.line.fill.background()

case_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(1.5))
case_frame = case_box.text_frame
case_frame.word_wrap = True

case_header = case_frame.paragraphs[0]
case_header.text = "典型案例"
case_header.font.size = Pt(18)
case_header.font.bold = True
case_header.font.color.rgb = COLOR_ACCENT
case_header.space_after = Pt(8)

case_content = case_frame.add_paragraph()
case_content.text = "招聘'海外兼职打字员'，通过Telegram联系，要求缴纳500元'会员费'才能接单。付费后只接到少量任务，收益远低于承诺，无法提现。"
case_content.font.size = Pt(13)
case_content.font.color.rgb = COLOR_TEXT

feature_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(1.5))
feature_frame = feature_box.text_frame
feature_frame.word_wrap = True

feature_header = feature_frame.paragraphs[0]
feature_header.text = "识别特征"
feature_header.font.size = Pt(18)
feature_header.font.bold = True
feature_header.font.color.rgb = COLOR_WARNING
feature_header.space_after = Pt(8)

features = ["通过非正规渠道联系", "要求缴纳会员费", "收益承诺过高", "提现条件苛刻"]
for f in features:
    p = feature_frame.add_paragraph()
    p.text = f"• {f}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

avoid_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.8))
avoid_frame = avoid_box.text_frame
avoid_frame.word_wrap = True

avoid_header = avoid_frame.paragraphs[0]
avoid_header.text = "规避方法"
avoid_header.font.size = Pt(18)
avoid_header.font.bold = True
avoid_header.font.color.rgb = COLOR_SUCCESS
avoid_header.space_after = Pt(8)

avoids = [
    "优先选择正规招聘平台", "警惕私密通讯工具：要求用Telegram等联系需警惕",
    "不预付任何费用", "核实公司真实性：跨境工作尤其要仔细核实"
]

for a in avoids:
    p = avoid_frame.add_paragraph()
    p.text = f"✓ {a}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

# ========== 陷阱6: 实习转正骗局 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "06 实习转正骗局"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(28)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_ACCENT

bg_box = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.15), Inches(0.5))
bg_box.fill.solid()
bg_box.fill.fore_color.rgb = COLOR_ACCENT
bg_box.line.fill.background()

case_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(1.5))
case_frame = case_box.text_frame
case_frame.word_wrap = True

case_header = case_frame.paragraphs[0]
case_header.text = "典型案例"
case_header.font.size = Pt(18)
case_header.font.bold = True
case_header.font.color.rgb = COLOR_ACCENT
case_header.space_after = Pt(8)

case_content = case_frame.add_paragraph()
case_content.text = "招聘实习生，承诺'表现优秀可转正'。实习6个月后，公司以各种理由拒绝转正，继续招聘新实习生，本质是利用廉价劳动力。"
case_content.font.size = Pt(13)
case_content.font.color.rgb = COLOR_TEXT

feature_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(1.5))
feature_frame = feature_box.text_frame
feature_frame.word_wrap = True

feature_header = feature_frame.paragraphs[0]
feature_header.text = "识别特征"
feature_header.font.size = Pt(18)
feature_header.font.bold = True
feature_header.font.color.rgb = COLOR_WARNING
feature_header.space_after = Pt(8)

features = ["转正承诺模糊", "实习期过长", "大量招聘实习生", "实习生离职率高"]
for f in features:
    p = feature_frame.add_paragraph()
    p.text = f"• {f}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

avoid_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.8))
avoid_frame = avoid_box.text_frame
avoid_frame.word_wrap = True

avoid_header = avoid_frame.paragraphs[0]
avoid_header.text = "规避方法"
avoid_header.font.size = Pt(18)
avoid_header.font.bold = True
avoid_header.font.color.rgb = COLOR_SUCCESS
avoid_header.space_after = Pt(8)

avoids = [
    "签订实习协议：明确转正条件和时间", "要求书面承诺：口头承诺无效",
    "设定实习期限：一般不超过6个月", "准备备选方案：不要过度依赖单一转正机会"
]

for a in avoids:
    p = avoid_frame.add_paragraph()
    p.text = f"✓ {a}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

# ========== 陷阱7: 信息泄露 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "07 信息泄露陷阱"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(28)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_ACCENT

bg_box = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.15), Inches(0.5))
bg_box.fill.solid()
bg_box.fill.fore_color.rgb = COLOR_ACCENT
bg_box.line.fill.background()

case_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(1.5))
case_frame = case_box.text_frame
case_frame.word_wrap = True

case_header = case_frame.paragraphs[0]
case_header.text = "典型案例"
case_header.font.size = Pt(18)
case_header.font.bold = True
case_header.font.color.rgb = COLOR_ACCENT
case_header.space_after = Pt(8)

case_content = case_frame.add_paragraph()
case_content.text = "求职者投递简历后，接到大量骚扰电话、营销短信。原来该公司将收集到的简历信息出售给培训机构、贷款公司等第三方。"
case_content.font.size = Pt(13)
case_content.font.color.rgb = COLOR_TEXT

feature_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(1.5))
feature_frame = feature_box.text_frame
feature_frame.word_wrap = True

feature_header = feature_frame.paragraphs[0]
feature_header.text = "识别特征"
feature_header.font.size = Pt(18)
feature_header.font.bold = True
feature_header.font.color.rgb = COLOR_WARNING
feature_header.space_after = Pt(8)

features = ["投递后频繁被骚扰", "要求提供过多个人信息", "公司知名度低", "招聘信息粗糙"]
for f in features:
    p = feature_frame.add_paragraph()
    p.text = f"• {f}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

avoid_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.8))
avoid_frame = avoid_box.text_frame
avoid_frame.word_wrap = True

avoid_header = avoid_frame.paragraphs[0]
avoid_header.text = "规避方法"
avoid_header.font.size = Pt(18)
avoid_header.font.bold = True
avoid_header.font.color.rgb = COLOR_SUCCESS
avoid_header.space_after = Pt(8)

avoids = [
    "优先选择正规平台：智联、前程无忧等", "谨慎提供敏感信息：身份证号、银行卡号不轻易给",
    "使用专用求职邮箱：避免主邮箱被骚扰", "发现泄露及时举报：向平台投诉"
]

for a in avoids:
    p = avoid_frame.add_paragraph()
    p.text = f"✓ {a}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

# ========== 陷阱8: 加班文化 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "08 加班文化隐性陷阱"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(28)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_ACCENT

bg_box = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.15), Inches(0.5))
bg_box.fill.solid()
bg_box.fill.fore_color.rgb = COLOR_ACCENT
bg_box.line.fill.background()

case_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(1.5))
case_frame = case_box.text_frame
case_frame.word_wrap = True

case_header = case_frame.paragraphs[0]
case_header.text = "典型案例"
case_header.font.size = Pt(18)
case_header.font.bold = True
case_header.font.color.rgb = COLOR_ACCENT
case_header.space_after = Pt(8)

case_content = case_frame.add_paragraph()
case_content.text = "公司表面上实行'朝九晚五双休'，实际要求员工经常'自愿加班'，加班费按最低标准甚至不发。入职后才发现工作强度极大，长期996。"
case_content.font.size = Pt(13)
case_content.font.color.rgb = COLOR_TEXT

feature_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(1.5))
feature_frame = feature_box.text_frame
feature_frame.word_wrap = True

feature_header = feature_frame.paragraphs[0]
feature_header.text = "识别特征"
feature_header.font.size = Pt(18)
feature_header.font.bold = True
feature_header.font.color.rgb = COLOR_WARNING
feature_header.space_after = Pt(8)

features = ["面试时暗示加班", "加班费计算不明确", "员工离职率高", "公司氛围压抑"]
for f in features:
    p = feature_frame.add_paragraph()
    p.text = f"• {f}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

avoid_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.8))
avoid_frame = avoid_box.text_frame
avoid_frame.word_wrap = True

avoid_header = avoid_frame.paragraphs[0]
avoid_header.text = "规避方法"
avoid_header.font.size = Pt(18)
avoid_header.font.bold = True
avoid_header.font.color.rgb = COLOR_SUCCESS
avoid_header.space_after = Pt(8)

avoids = [
    "面试时确认工作制度：明确加班政策", "查询员工评价：脉脉、知乎查看真实反馈",
    "观察办公环境：面试时留意员工状态", "设定个人底线：明确能接受的加班范围"
]

for a in avoids:
    p = avoid_frame.add_paragraph()
    p.text = f"✓ {a}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(4)

# ========== 总结页 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
title_frame = title.text_frame
title_frame.text = "总结"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(32)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_PRIMARY

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(2.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

summary = [
    "求职陷阱无处不在，但并非无法防范",
    "",
    "核心原则：",
    "• 不轻信：对过高承诺保持警惕",
    "• 不付费：正规招聘不收取任何费用",
    "• 多核实：通过多种渠道验证公司信息",
    "• 保留证据：合同、付款凭证、聊天记录",
    "",
    "遇到问题时，及时向劳动监察部门、网警平台举报"
]

for line in summary:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT
    if "•" in line:
        p.space_before = Pt(6)
        p.space_after = Pt(3)
    elif line == "":
        p.space_after = Pt(12)
    else:
        p.space_after = Pt(6)

# ========== 结束页 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])

end_text = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(1))
end_frame = end_text.text_frame
end_frame.text = "祝你求职顺利！"
end_p = end_frame.paragraphs[0]
end_p.font.size = Pt(36)
end_p.font.bold = True
end_p.font.color.rgb = COLOR_PRIMARY
end_p.alignment = PP_ALIGN.CENTER

footer = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.5))
footer_frame = footer.text_frame
footer_frame.text = "保持警惕，理性求职"
footer_p = footer_frame.paragraphs[0]
footer_p.font.size = Pt(18)
footer_p.font.color.rgb = COLOR_SECONDARY
footer_p.alignment = PP_ALIGN.CENTER

# 保存文件
prs.save('求职常见陷阱类型.pptx')
print("PPT generated successfully: 求职常见陷阱类型.pptx")
