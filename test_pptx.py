from pptx import Presentation
from pptx.util import Inches, Pt

# 创建新的演示文稿
prs = Presentation()

# 设置幻灯片尺寸为16:9
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# 添加标题幻灯片
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)

# 设置标题
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, yy!"
subtitle.text = "这是一个测试 PPT - 由 python-pptx 生成"

# 添加内容幻灯片
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "我能做什么"

tf = body_shape.text_frame
tf.text = "使用 python-pptx 库"

p = tf.add_paragraph()
p.text = "创建精美的 PowerPoint 演示文稿"
p.level = 1

p = tf.add_paragraph()
p.text = "支持文本、图片、图表等内容"
p.level = 1

p = tf.add_paragraph()
p.text = "可以自定义布局和样式"
p.level = 1

# 保存文件
prs.save('test_presentation.pptx')
print("Test PPT generated: test_presentation.pptx")
